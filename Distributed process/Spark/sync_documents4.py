#!/usr/bin/env python3

import argparse, io, os, re
from pyspark.sql import SparkSession, Row
from pyspark.sql.functions import col, regexp_extract, input_file_name, lit
from pyspark.sql.types import StringType, IntegerType
from pyspark.sql.types import StructType, StructField, StringType
from pyspark.sql.functions import col as Fcol, length, trim

import re

_WHITES_RE = re.compile(r"\s+", re.UNICODE)
_NON_PRINTABLE_RE = re.compile(r"[^\x20-\x7E\u00A0-\uFFFF]", re.UNICODE)

def normalize_text(s: str) -> str:
    if not s:
        return ""
    # supprime les octets/contrôles non imprimables
    s = _NON_PRINTABLE_RE.sub("", s)
    # remplace tout whitespace (espaces, tabs, \r, \n, NBSP…) par un simple espace
    s = _WHITES_RE.sub(" ", s)
    return s.strip()


# ---------- extractors ----------
def extract_pdf(bytes_data, min_chars=200):
    try:
        from pdfminer.high_level import extract_text
        text = extract_text(io.BytesIO(bytes_data)) or ""
        text = text.strip()
        return text if len(text) >= min_chars else ""
    except Exception:
        return ""

def extract_docx(bytes_data, min_chars=200):
    try:
        import docx
        bio = io.BytesIO(bytes_data)
        doc = docx.Document(bio)
        text = "\n".join(p.text for p in doc.paragraphs).strip()
        return text if len(text) >= min_chars else ""
    except Exception:
        return ""

def extract_pptx(bytes_data, min_chars=200, max_slides=60):
    try:
        from pptx import Presentation
        bio = io.BytesIO(bytes_data)
        prs = Presentation(bio)
        parts = []
        # NE PAS dépasser max_slides
        for i, slide in enumerate(prs.slides, 1):
            if i > max_slides:
                break
            buf = []
            for shp in slide.shapes:
                if hasattr(shp, "text"):
                    buf.append(shp.text)
            if buf:
                parts.append(f"[SLIDE {i}]\n" + "\n".join(buf))
        text = "\n\n".join(parts).strip()
        return text if len(text) >= min_chars else ""
    except Exception:
        return ""

def extract_xlsx(bytes_data, sample_rows=300, min_chars=200):
    try:
        import openpyxl
        bio = io.BytesIO(bytes_data)
        wb = openpyxl.load_workbook(bio, data_only=True, read_only=True)
        parts = []
        for ws in wb.worksheets:
            rows = ws.iter_rows(values_only=True)
            head = next(rows, None)
            if head:
                head = [str(c) if c is not None else "" for c in head]
            buf = []
            buf.append(f"[SHEET] {ws.title}")
            if head:
                buf.append("[HEAD] " + " | ".join(head))
            for idx, r in enumerate(rows, 1):
                if idx > sample_rows: break
                row_txt = " | ".join("" if v is None else str(v) for v in r)
                buf.append(row_txt)
            parts.append("\n".join(buf))
        text = "\n\n".join(parts).strip()
        return text if len(text) >= min_chars else ""
    except Exception:
        return ""

def extract_csv(bytes_data, sample_rows=300, min_chars=200):
    try:
        import csv
        bio = io.BytesIO(bytes_data)
        text = bio.read().decode("utf-8", errors="ignore")
        lines = text.splitlines()
        if not lines:
            return ""
        reader = csv.reader(lines)
        out = []
        for i, row in enumerate(reader):
            if i == 0:
                out.append("[HEAD] " + " | ".join(row))
            else:
                out.append(" | ".join(row))
            if i >= sample_rows: break
        txt = "\n".join(out).strip()
        return txt if len(txt) >= min_chars else ""
    except Exception:
        return ""

def extract_xls(bytes_data, sample_rows=300, min_chars=200):
    try:
        import xlrd  # xlrd>=2.0.1 recommandé
        wb = xlrd.open_workbook(file_contents=bytes_data)  # bytes, pas bytearray
        parts = []
        for s in wb.sheets():
            buf = [f"[SHEET] {s.name}"]
            if s.nrows > 0:
                head = [str(s.cell_value(0, c)) for c in range(s.ncols)]
                buf.append("[HEAD] " + " | ".join(head))
            for r in range(1, min(s.nrows, sample_rows + 1)):
                row = [str(s.cell_value(r, c)) for c in range(s.ncols)]
                buf.append(" | ".join(row))
            parts.append("\n".join(buf))
        text = "\n\n".join(parts).strip()
        return text if len(text) >= min_chars else ""
    except Exception as e:
        return ""


def extract_xml(bytes_data, min_chars=100):
    try:
        import xml.etree.ElementTree as ET
        root = ET.fromstring(bytes_data.decode("utf-8", errors="ignore"))
        texts = []
        for elem in root.iter():
            if elem.text and elem.text.strip():
                texts.append(elem.text.strip())
        text = "\n".join(texts)
        return text if len(text) >= min_chars else ""
    except Exception:
        return ""
    
def ocr_image(bytes_data, min_chars=120, tesseract_cmd="/usr/bin/tesseract", lang="eng+fra"):
    try:
        from PIL import Image
        import pytesseract, io
        if tesseract_cmd:
            pytesseract.pytesseract.tesseract_cmd = tesseract_cmd
        img = Image.open(io.BytesIO(bytes_data))
        text = pytesseract.image_to_string(img, lang=lang).strip()
        return text if len(text) >= min_chars else ""
    except Exception:
        return ""



def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--in", dest="in_dir", required=True)
    ap.add_argument("--out", dest="out_dir", required=True)
    ap.add_argument("--tmp", dest="tmp_dir", required=True)
    ap.add_argument("--sample_rows", type=int, default=300)
    ap.add_argument("--min_chars", type=int, default=200)
    ap.add_argument("--enable_ocr", type=str, default="false")  
    ap.add_argument("--tesseract_cmd", type=str, default="") 
    ap.add_argument("--ocr_lang", type=str, default="eng+fra")    
    ap.add_argument("--max_input_bytes", type=int, default=25*1024*1024)  
    ap.add_argument("--max_output_chars", type=int, default=200_000)    
    ap.add_argument("--max_slides", type=int, default=60) 
    args = ap.parse_args()

    spark = SparkSession.builder.appName("ExtractDocsToTxt").enableHiveSupport().getOrCreate()
    sc = spark.sparkContext

    # 1) Lire les binaires depuis HDFS (miroir)
    # DataFrame colonnes: path, modificationTime, length, content, ...
    df_bin = spark.read.format("binaryFile") \
        .option("pathGlobFilter", "*") \
        .load(args.in_dir)
    
    df_bin = df_bin.filter(col("length") <= args.max_input_bytes)

    # 2) Extraire id et extension depuis le nom de fichier {id}.{ext}
    #    Exemple path: hdfs:///data/docs_bin/4a3ba7....docx
    df = df_bin.withColumn("filename", regexp_extract(input_file_name(), r"([^/]+)$", 1))
    df = df.withColumn("id", regexp_extract(col("filename"), r"^([^.]+)", 1)) \
           .withColumn("ext", regexp_extract(col("filename"), r"\.([^.]+)$", 1))

    # 3) Transformer en RDD et appliquer l'extraction par extension
    def process_partition(rows_iter):
        import json
        for r in rows_iter:
            res = {"id": r.id, "ext": r.ext.lower(), "status": "skip", "reason": ""}
            try:
                ext = (r.ext or "").lower()
                content = r.content
                content_bytes = bytes(r.content)
                _id = r.id
                text = ""
                if ext == "pdf":
                    text = extract_pdf(content, min_chars=args.min_chars)
                elif ext == "docx":
                    text = extract_docx(content, min_chars=args.min_chars)
                elif ext == "pptx":
                    text = extract_pptx(content, min_chars=args.min_chars, max_slides=args.max_slides)

                elif ext in ("xlsx", "xlsm"):
                    text = extract_xlsx(content, sample_rows=args.sample_rows, min_chars=args.min_chars)
                elif ext == "csv":
                    text = extract_csv(content, sample_rows=args.sample_rows, min_chars=args.min_chars)
                elif ext == "xls":
                    # text = extract_xls(content, sample_rows=args.sample_rows, min_chars=args.min_chars)
                    text = extract_xls(content_bytes, sample_rows=args.sample_rows, min_chars=args.min_chars)
                elif ext == "xml":
                    text = extract_xml(content, min_chars=max(100, args.min_chars // 2))
                elif ext in ("jpg", "jpeg", "png"):
                    if args.enable_ocr.lower() == "true":
                        text = ocr_image(
                            content_bytes,
                            min_chars=max(80, args.min_chars // 2),
                            tesseract_cmd=args.tesseract_cmd or "/usr/bin/tesseract",
                            lang=args.ocr_lang
                        )
                elif ext in ("zip", "qgz", "gsheet", "gslides"):
                    res["status"] = "skip"
                    res["reason"] = f"unsupported_ext_{ext}"
                    yield Row(id=r.id, kind="audit", text=None, detail=json.dumps(res))
                    continue 

                text = normalize_text(text)

                min_chars = max(1, args.min_chars)  
                if text and len(text) > args.max_output_chars:
                    text = text[:args.max_output_chars]

                if text and len(text) >= min_chars:
                    res["status"] = "ok"
                    yield Row(id=r.id, kind="txt", text=text, detail=None)
                else:
                    res["status"] = "empty"
                    res["reason"] = "no_text_or_below_min_chars"
                    yield Row(id=r.id, kind="audit", text=None, detail=json.dumps(res))
            except Exception as e:
                res["status"] = "error"
                res["reason"] = repr(e)
                yield Row(kind="audit", id=r.id, detail=json.dumps(res))

    rdd_txt = df.select("id", "ext", "content").rdd.mapPartitions(process_partition)

    schema = StructType([
    StructField("id", StringType(), False),
    StructField("kind", StringType(), False),    # "txt" ou "audit"
    StructField("text", StringType(), True),     # présent si kind == "txt"
    StructField("detail", StringType(), True),   # présent si kind == "audit"
])

    df_txt = spark.createDataFrame(rdd_txt, schema=schema)
    # df_txt = spark.createDataFrame(rdd_txt) 

    # 4) Ecrire en TEXT partitionné par id → .../id=<id>/part-0000...
    #    (format TEXT attend une colonne "value")
    from pyspark.sql.functions import col as Fcol
    # df_txt.write.mode("overwrite").format("text").partitionBy("id").save(args.tmp_dir)
    
    min_chars_lit = int(args.min_chars)

    df_ok = (
        df_txt
        .filter(Fcol("kind") == "txt")
        .withColumn("value", trim(Fcol("text")))
        .filter(Fcol("value").isNotNull())
        .filter(length(Fcol("value")) >= min_chars_lit)
        .select(Fcol("id"), Fcol("value"))
    )

    if df_ok.rdd.isEmpty():
        print("Aucun document texte à écrire (tous vides/sous seuil/skippés).")
    else:
        df_ok.repartition("id") \
            .write.mode("overwrite").format("text") \
            .partitionBy("id").save(args.tmp_dir)

    # 5) Renommer chaque part-*.txt en {id}.txt dans args.out (driver, API Hadoop)
    jvm = spark._jvm
    conf = spark._jsc.hadoopConfiguration()
    fs = jvm.org.apache.hadoop.fs.FileSystem.get(conf)
    Path = jvm.org.apache.hadoop.fs.Path

    fs.mkdirs(Path(args.out_dir))

    for status in fs.listStatus(Path(args.tmp_dir)):
        if status.isDirectory():
            part_dir = status.getPath()               # .../id=<id>
            name = part_dir.getName()                 # id=<id>
            if not name.startswith("id="): 
                continue
            doc_id = name.split("=", 1)[1]
            files = fs.listStatus(part_dir)
            # prendre le premier part-*.txt non vide
            part_file = None
            for f in files:
                if f.isFile() and f.getPath().getName().startswith("part-"):
                    part_file = f.getPath()
                    break
            if part_file:
                target = Path(args.out_dir + "/" + doc_id + ".txt")
                # s'il existe déjà, supprimer/overwrite
                if fs.exists(target):
                    fs.delete(target, False)
                fs.rename(part_file, target)

    # Nettoyage du répertoire tmp si tu veux :
    fs.delete(Path(args.tmp_dir), True)

    spark.stop()

if __name__ == "__main__":
    main()
